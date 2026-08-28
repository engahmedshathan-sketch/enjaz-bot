import os
import io
import asyncio
import uuid
import sys
import requests
import random
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import (
    Application, CommandHandler, CallbackQueryHandler,
    MessageHandler, filters, ContextTypes
)
from aiohttp import web

try:
    import pypdf
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

TELEGRAM_TOKEN = os.environ.get("TELEGRAM_TOKEN", "8867458917:AAFVTZ4NJLzTFdon-Z0Zf--oyRxt2u208JI")

GRADES = {
    "kg": "رياض الأطفال", "p1": "الأول الابتدائي", "p2": "الثاني الابتدائي", "p3": "الثالث الابتدائي",
    "p4": "الرابع الابتدائي", "p5": "الخامس الابتدائي", "p6": "السادس الابتدائي",
    "m1": "الأول المتوسط", "m2": "الثاني المتوسط", "m3": "الثالث المتوسط",
    "s1": "الأول الثانوي", "s2": "الثاني الثانوي", "s3": "الثالث الثانوي",
}

SYSTEM_PROMPT = """
أنت خبير أكاديمي ومحكم علمي في الدراسات العليا وإدارة الأعمال وباحث تربوي في مقررات وزارة التعليم بالجمهورية اليمنية والمملكة العربية السعودية لعام 1448هـ.
التعليمات الصارمة:
1. عند طلب البحوث الأكاديمية: اكتب بحوثاً وخططاً أكاديمية شاملة، مفصلة، وعميقة جداً على غرار رسائل الماجستير في جامعة صعدة (مقدمة، مشكلة، أسئلة، فرضيات، أهداف، أهمية، حدود، مصطلحات، دراسات سابقة مع التعليق، إطار منهجي، ومراجع).
2. عند طلب الخدمات المدرسية: التزم بنصوص المنهج المدرسي الرسمي السعودي واليمني والمصطلحات الدقيقة.
"""

COLOR_PALETTES = [
    {"primary": RGBColor(27, 73, 101), "accent": RGBColor(95, 168, 211), "card": RGBColor(248, 250, 252)},
    {"primary": RGBColor(44, 122, 123), "accent": RGBColor(129, 230, 217), "card": RGBColor(240, 253, 250)},
    {"primary": RGBColor(88, 28, 135), "accent": RGBColor(196, 181, 253), "card": RGBColor(245, 243, 255)},
    {"primary": RGBColor(180, 83, 9), "accent": RGBColor(252, 211, 77), "card": RGBColor(255, 251, 235)},
    {"primary": RGBColor(159, 18, 57), "accent": RGBColor(253, 164, 175), "card": RGBColor(255, 241, 242)},
    {"primary": RGBColor(6, 95, 70), "accent": RGBColor(110, 231, 183), "card": RGBColor(236, 253, 245)},
]

def query_ai_engine(prompt: str) -> str:
    unique_id = str(uuid.uuid4())
    cache_buster_prompt = f"{prompt}\n\n[معرف فريد: {unique_id} - التزم بالمنهجية الأكاديمية والمدرسية الشاملة لعام 1448هـ]"
    payloads = [
        {
            "url": f"https://text.pollinations.ai/?seed={random.randint(1, 999999)}",
            "data": {
                "messages": [
                    {"role": "system", "content": SYSTEM_PROMPT},
                    {"role": "user", "content": cache_buster_prompt},
                ],
                "model": "openai",
                "jsonMode": False
            },
        },
        {
            "url": "https://api.airforce/v1/chat/completions",
            "data": {
                "model": "gpt-4o-mini",
                "messages": [
                    {"role": "system", "content": SYSTEM_PROMPT},
                    {"role": "user", "content": cache_buster_prompt},
                ],
            },
        },
    ]
    for payload in payloads:
        try:
            response = requests.post(payload["url"], json=payload["data"], timeout=90)
            if response.status_code == 200:
                text = response.text.strip()
                if text:
                    try:
                        data = response.json()
                        if "choices" in data:
                            return data["choices"][0]["message"]["content"].strip()
                    except:
                        pass
                    if len(text) > 200:
                        return text
        except:
            continue
    return "محتوى تعليمي معتمد لمنصة إنجاز 1448هـ."

def clean_pdf_text_with_ai(raw_text: str) -> str:
    prompt = f"""
قم بتنظيف وترتيب النص التالي المستخرج من ملف PDF رسمي ليكون بصيغة تقرير أو مستند منظم تماماً باللغة العربية الفصحى بدون أي أحرف مقلوبة أو مبعثرة:
{raw_text[:4000]}
"""
    try:
        response = requests.post(
            "https://text.pollinations.ai/",
            json={"messages": [{"role": "user", "content": prompt}], "model": "openai", "jsonMode": False},
            timeout=30
        )
        if response.status_code == 200 and len(response.text) > 50:
            return response.text.strip()
    except:
        pass
    return raw_text

def extract_pdf_to_text(pdf_bytes: bytes) -> str:
    if not PDF_SUPPORT:
        return "مكتبة قراءة الـ PDF غير متوفرة."
    try:
        reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
        full_text = ""
        for page in reader.pages:
            t = page.extract_text()
            if t:
                full_text += t + "\n"
        return full_text.strip()
    except Exception as e:
        return f"خطأ: {str(e)}"

def fetch_unique_slide_image(slide_index: int, topic: str) -> io.BytesIO:
    keywords = ["school,students", "classroom,learning", "science,experiment", "math,numbers", "library,books"]
    kw = keywords[(slide_index - 1) % len(keywords)]
    lock = random.randint(1, 999999) 
    url = f"https://loremflickr.com/600/450/{kw}?lock={lock}"
    try:
        response = requests.get(url, timeout=6)
        if response.status_code == 200 and len(response.content) > 3000:
            stream = io.BytesIO(response.content)
            stream.seek(0)
            return stream
    except:
        pass
    img = Image.new("RGB", (600, 450), color="#F8FAFC")
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle([15, 15, 585, 435], radius=15, fill="#FFFFFF", outline="#CBD5E1", width=2)
    output = io.BytesIO()
    img.save(output, format="PNG")
    output.seek(0)
    return output

def generate_dynamic_30_slides_data(grade: str, subject: str, topic: str):
    grade = (grade or "غير محدد").strip()
    subject = (subject or "غير محددة").strip()
    topic = (topic or "الدرس").strip()

    ai_prompt = f"""
أنت تقوم بتفريغ محتوى المنهج التعليمي لعمل عرض بوربوينت لدرس محدد.
الصف: {grade}
المادة: {subject}
الموضوع/الدرس: {topic}
مطلوب منك 30 شريحة تعليمية احترافية.
التنسيق الإلزامي لكل شريحة:
---SLIDE---
TITLE: [اكتب عنواناً فرعياً من داخل الدرس]
CONTENT:
- [نص أو تعريف أساسي]
- [مثال أو تمرين تطبيقي]
- [سؤال تقويمي للطلاب]
"""
    ai_raw = query_ai_engine(ai_prompt)
    ai_slides = []
    if ai_raw:
        chunks = ai_raw.split("---SLIDE---")
        for chunk in chunks:
            if not chunk.strip(): continue
            lines = [l.strip() for l in chunk.splitlines() if l.strip()]
            title = ""
            points = []
            for line in lines:
                if line.startswith("TITLE:"):
                    title = line.replace("TITLE:", "").strip()
                elif line.startswith("-") or line.startswith("•"):
                    points.append(line.lstrip("-•* ").strip())
            if title and points:
                ai_slides.append((title, points[:5]))

    while len(ai_slides) < 30:
        idx = len(ai_slides) + 1
        ai_slides.append((f"شريحة تعليمية ({idx}) - {topic}", [
            f"مفهوم أساسي في مادة {subject} للدرس {topic}.",
            "تطبيق عملي وتدريب للطلاب.",
            "سؤال تقويمي استدلالي."
        ]))

    return ai_slides[:30]

def create_powerpoint_presentation_full(grade: str, subject: str, topic: str, output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slides_data = generate_dynamic_30_slides_data(grade, subject, topic)

    for idx, (title_text, points) in enumerate(slides_data, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        palette = COLOR_PALETTES[(idx - 1) % len(COLOR_PALETTES)]
        layout_style = idx % 3  

        if layout_style == 0:
            top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.25))
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = palette["primary"]
            top_bar.line.fill.background()

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.5), Inches(11.333), Inches(5.2))
            card.fill.solid()
            card.fill.fore_color.rgb = palette["card"]
            card.line.color.rgb = palette["accent"]
            card.line.width = Pt(1.5)

            title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(10.333), Inches(0.8))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = title_text
            p_title.alignment = PP_ALIGN.RIGHT
            p_title.font.size = Pt(24)
            p_title.font.bold = True
            p_title.font.color.rgb = palette["primary"]

            content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(4.5))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True

            for p_idx, point in enumerate(points):
                p = tf_content.paragraphs[0] if p_idx == 0 else tf_content.add_paragraph()
                p.text = f"🔹 {point}"
                p.alignment = PP_ALIGN.RIGHT
                p.font.size = Pt(17)
                p.font.color.rgb = RGBColor(30, 41, 59)
                p.space_after = Pt(14)

        elif layout_style == 1:
            side_banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
            side_banner.fill.solid()
            side_banner.fill.fore_color.rgb = palette["accent"]
            side_banner.line.fill.background()

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.4), Inches(7.4), Inches(5.3))
            card.fill.solid()
            card.fill.fore_color.rgb = palette["card"]
            card.line.color.rgb = palette["primary"]
            card.line.width = Pt(1.5)

            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = title_text
            p_title.alignment = PP_ALIGN.RIGHT
            p_title.font.size = Pt(24)
            p_title.font.bold = True
            p_title.font.color.rgb = palette["primary"]

            content_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.6), Inches(6.8), Inches(4.8))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True

            for p_idx, point in enumerate(points):
                p = tf_content.paragraphs[0] if p_idx == 0 else tf_content.add_paragraph()
                p.text = f"◀ {point}"
                p.alignment = PP_ALIGN.RIGHT
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(30, 41, 59)
                p.space_after = Pt(12)

            img_stream = fetch_unique_slide_image(idx, topic)
            if img_stream:
                slide.shapes.add_picture(img_stream, Inches(0.8), Inches(1.4), Inches(4.1), Inches(5.3))

        else:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(7.4), Inches(5.3))
            card.fill.solid()
            card.fill.fore_color.rgb = palette["card"]
            card.line.color.rgb = palette["accent"]
            card.line.width = Pt(1.5)

            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = title_text
            p_title.alignment = PP_ALIGN.RIGHT
            p_title.font.size = Pt(24)
            p_title.font.bold = True
            p_title.font.color.rgb = palette["primary"]

            content_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.6), Inches(6.8), Inches(4.8))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True

            for p_idx, point in enumerate(points):
                p = tf_content.paragraphs[0] if p_idx == 0 else tf_content.add_paragraph()
                p.text = f"• {point}"
                p.alignment = PP_ALIGN.RIGHT
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(30, 41, 59)
                p.space_after = Pt(12)

            img_stream = fetch_unique_slide_image(idx, topic)
            if img_stream:
                slide.shapes.add_picture(img_stream, Inches(8.4), Inches(1.4), Inches(4.1), Inches(5.3))

        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.4))
        p_foot = footer_box.text_frame.paragraphs[0]
        p_foot.text = f"شريحة {idx} من 30 | منصة إنجاز | {grade} | {subject} | 1448هـ"
        p_foot.alignment = PP_ALIGN.LEFT
        p_foot.font.size = Pt(10)
        p_foot.font.color.rgb = RGBColor(148, 163, 184)

    prs.save(output_path)

def create_ppt_from_pdf_text(pdf_text: str, output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    paragraphs = [p.strip() for p in pdf_text.split("\n") if p.strip() and len(p.strip()) > 5]
    if not paragraphs: paragraphs = ["محتوى مستخرج من ملف الـ PDF"]

    chunk_size = 4
    chunks = [paragraphs[i:i + chunk_size] for i in range(0, len(paragraphs), chunk_size)]
    total_slides = min(max(len(chunks), 1), 30)

    for idx in range(1, total_slides + 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        palette = COLOR_PALETTES[(idx - 1) % len(COLOR_PALETTES)]
        
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.25))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = palette["primary"]
        top_bar.line.fill.background()
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(1.5), Inches(7.2), Inches(5.1))
        card.fill.solid()
        card.fill.fore_color.rgb = palette["card"]
        card.line.color.rgb = palette["accent"]
        card.line.width = Pt(1.5)

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = f"شريحة ملف الإنجاز رقم ({idx})"
        p_title.alignment = PP_ALIGN.RIGHT
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = palette["primary"]

        content_box = slide.shapes.add_textbox(Inches(5.6), Inches(1.7), Inches(6.8), Inches(4.7))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True

        current_chunk = chunks[(idx - 1) % len(chunks)]
        for p_idx, point in enumerate(current_chunk):
            p = tf_content.paragraphs[0] if p_idx == 0 else tf_content.add_paragraph()
            p.text = f"◀ {point[:110]}"
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(30, 41, 59)
            p.space_after = Pt(12)

        img_stream = fetch_unique_slide_image(idx, "pdf_import")
        if img_stream:
            slide.shapes.add_picture(img_stream, Inches(0.8), Inches(1.5), Inches(4.3), Inches(5.1))

        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.4))
        p_foot = footer_box.text_frame.paragraphs[0]
        p_foot.text = f"شريحة {idx} من {total_slides} | منصة إنجاز | تحويل PDF مباشر | 1448هـ"
        p_foot.alignment = PP_ALIGN.LEFT
        p_foot.font.size = Pt(11)
        p_foot.font.color.rgb = RGBColor(148, 163, 184)

    prs.save(output_path)

def create_educational_doc_1448(service_code: str, grade: str, subject: str, topic: str, output_path: str):
    doc = Document()
    for section in doc.sections:
        section.top_margin, section.bottom_margin = DocxInches(1), DocxInches(1)
        section.left_margin, section.right_margin = DocxInches(1), DocxInches(1)

    if service_code == "svc_research":
        prompt = f"""
قم بإعداد خطة وبحث أكاديمي جامعياً مفصلاً، شاملاً، وعميقاً جداً حول الموضوع التالي: '{topic}'.
يجب أن يكون البحث على غرار رسائل الماجستير والدكتوراه الأكاديمية في الجامعات اليمنية (مثل جامعة صعدة) ويحتوي حصرياً على الأقسام التالية مفصلة بالكامل:
1. العنوان الرئيسي بالعربية والإنجليزية وبيانات الإعداد والإشراف والجامعة (جامعة صعدة - نيابة الدراسات العليا).
2. أولاً: المقدمة (مقدمة مفصلة وعميقة تبين الأهمية العلمية والعملية للموضوع وتأصيلها).
3. ثانياً: مشكلة الدراسة وأسئلتها (تحديد دقيق للمشكلة وصياغة السؤال الرئيس والأسئلة الفرعية).
4. ثالثاً: فرضيات الدراسة (الفرضيات الرئيسية والفرعية المرتبطة بالمتغيرات).
5. رابعاً: أهداف الدراسة (الهدف الرئيسي والأهداف الفرعية).
6. خامساً: أهمية الدراسة (الأهمية العلمية والأهمية التطبيقية).
7. سادساً: حدود الدراسة (الموضوعية، البشرية، المكانية مثل محافظة صعدة، والزمانية 1448هـ).
8. سابعاً: نموذج الدراسة ومتغيراتها (المتغير المستقل، التابع، الديموغرافية).
9. ثامناً: مصطلحات الدراسة (التعريف الاصطلاحي والإجرائي للمتغيرات).
10. تاسعاً: الدراسات السابقة (عرض 5 دراسات سابقة حديثة محلية وعربية وأجنبية مع التعليق عليها).
11. عاشراً: الإطار المنهجي والميداني (منهج الدراسة، المجتمع والعينة، الأداة، الأساليب الإحصائية).
12. قائمة المراجع (مراجع عربية وأجنبية موثقة).
اكتب محتوى غنياً، مفصلاً بالكامل، وبصياغة أكاديمية رصينة بدون أي اختصار.
"""
        doc_title = f"خطة وبحث أكاديمي متكامل\n{topic}\nجامعة صعدة - 1448هـ"
    else:
        prompts = {
            "svc_exam": f"اكتب اختباراً شاملاً لعام 1448هـ للصف {grade} مادة {subject} حول {topic} مع نموذج إجابة وتوزيع درجات.",
            "svc_remedial": f"اكتب خطة علاجية وإثرائية وأوراق عمل تفصيلية للصف {grade} مادة {subject} حول {topic}.",
            "svc_portfolio": f"اكتب ملف إنجاز إلكترونياً منظماً وشاملاً للمعلم للصف {grade} ومادة {subject} لعام 1448هـ.",
            "svc_performance": f"اكتب سجل ملف أداء وظيفي للمعلم لعام 1448هـ للصف {grade} ومادة {subject} مع بنود التقييم.",
            "svc_operation": f"اكتب خطة تشغيلية تعليمية متكاملة لعام 1448هـ للصف {grade} ومادة {subject} حول {topic}.",
            "svc_loss": f"اكتب خطة معالجة الفاقد التعليمي مفصلة للصف {grade} في مادة {subject} حول {topic}.",
        }
        titles = {
            "svc_exam": f"الاختبار وتحليل النتائج\n{subject} - {grade}\n{topic}",
            "svc_remedial": f"الخطة العلاجية والإثرائية\n{subject} - {grade}\n{topic}",
            "svc_portfolio": f"ملف الإنجاز الإلكتروني 1448هـ\n{subject} - {grade}",
            "svc_performance": f"ملف الأداء الوظيفي 1448هـ\n{subject} - {grade}",
            "svc_operation": f"الخطة التشغيلية 1448هـ\n{subject} - {grade}",
            "svc_loss": f"خطة معالجة الفاقد التعليمي 1448هـ\n{subject} - {grade}",
        }
        prompt = prompts.get(service_code, f"اكتب وثيقة تعليمية مفصلة لعام 1448هـ حول {topic}.")
        doc_title = titles.get(service_code, f"وثيقة تعليمية 1448هـ\n{topic}")

    ai_content = query_ai_engine(prompt)

    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_title = title_p.add_run(doc_title)
    run_title.font.size, run_title.font.bold, run_title.font.color.rgb = DocxPt(18), True, DocxRGB(27, 73, 101)

    sub_p = doc.add_paragraph()
    sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_sub = sub_p.add_run(f"منصة إنجاز | الصف: {grade} | المادة: {subject} | العام 1448هـ\n" + "—" * 40)
    run_sub.font.size, run_sub.font.color.rgb = DocxPt(11), DocxRGB(100, 116, 139)

    if len(ai_content) > 200:
        for block in ai_content.split("\n\n"):
            clean = block.strip()
            if not clean: continue
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            run = p.add_run(clean)
            run.font.size, run.font.color.rgb = DocxPt(11.5), DocxRGB(30, 41, 59)
            p.paragraph_format.line_spacing, p.paragraph_format.space_after = 1.25, DocxPt(6)
    else:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        p.add_run(f"تم إعداد المستند التعليمي للصف {grade} في مادة {subject} حول {topic} لعام 1448هـ.")

    doc.save(output_path)

def main_menu():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("📊 بوربوينت 30 شريحة ذكي بتصاميم متغيرة", callback_data="svc_ppt")],
        [InlineKeyboardButton("📝 اختبارات + جدول مواصفات + نافس", callback_data="svc_exam")],
        [InlineKeyboardButton("📈 خطط علاجية وإثرائية + أوراق عمل", callback_data="svc_remedial")],
        [InlineKeyboardButton("🗂 ملف إنجاز المعلم/المعلمة", callback_data="svc_portfolio")],
        [InlineKeyboardButton("📑 ملف الأداء الوظيفي", callback_data="svc_performance")],
        [InlineKeyboardButton("📅 الخطة التشغيلية", callback_data="svc_operation")],
        [InlineKeyboardButton("📚 خطة الفاقد التعليمي", callback_data="svc_loss")],
        [InlineKeyboardButton("🎓 بحث جامعي وأكاديمي شامل (جامعة صعدة)", callback_data="svc_research")],
        [InlineKeyboardButton("📊 تحويل ملف PDF إلى بوربوينت", callback_data="mode_pdf_ppt")],
        [InlineKeyboardButton("📝 تحويل وتعديل ملف PDF إلى Word منظم", callback_data="mode_pdf_word")],
        [InlineKeyboardButton("🎓 اختيار الصف الدراسي", callback_data="choose_grade")],
        [InlineKeyboardButton("🔄 تحديث وإعادة تشغيل البوت", callback_data="bot_restart")],
    ])

def grade_menu():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("رياض الأطفال", callback_data="grade_kg")],
        [InlineKeyboardButton("الأول ابتدائي", callback_data="grade_p1"), InlineKeyboardButton("الثاني ابتدائي", callback_data="grade_p2")],
        [InlineKeyboardButton("الثالث ابتدائي", callback_data="grade_p3"), InlineKeyboardButton("الرابع ابتدائي", callback_data="grade_p4")],
        [InlineKeyboardButton("الخامس ابتدائي", callback_data="grade_p5"), InlineKeyboardButton("السادس ابتدائي", callback_data="grade_p6")],
        [InlineKeyboardButton("الأول متوسط", callback_data="grade_m1"), InlineKeyboardButton("الثاني متوسط", callback_data="grade_m2")],
        [InlineKeyboardButton("الثالث متوسط", callback_data="grade_m3")],
        [InlineKeyboardButton("الأول ثانوي", callback_data="grade_s1"), InlineKeyboardButton("الثاني ثانوي", callback_data="grade_s2")],
        [InlineKeyboardButton("الثالث ثانوي", callback_data="grade_s3")],
        [InlineKeyboardButton("⬅️ القائمة الرئيسية", callback_data="home")],
    ])

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data.setdefault("grade", "")
    context.user_data.setdefault("subject", "")
    welcome_text = "🌟 *أهلاً بك في منصة إنجاز الشاملة والمحدثة 1448هـ*\n\n👇 اختر الخدمة المطلوبة:"
    if update.message:
        await update.message.reply_text(welcome_text, reply_markup=main_menu(), parse_mode="Markdown")
    elif update.callback_query:
        await update.callback_query.message.reply_text(welcome_text, reply_markup=main_menu(), parse_mode="Markdown")

async def restart_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    msg = update.message or update.callback_query.message
    await msg.reply_text("🔄 جاري إعادة تشغيل وتحديث البوت فوراً...")
    await asyncio.sleep(1)
    os.execl(sys.executable, sys.executable, *sys.argv)

async def button_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    data = query.data

    if data == "home":
        await query.edit_message_text("👇 القائمة الرئيسية:", reply_markup=main_menu())
        return

    if data == "choose_grade":
        await query.edit_message_text("🎓 اختر الصف الدراسي:", reply_markup=grade_menu())
        return

    if data.startswith("grade_"):
        grade_key = data.replace("grade_", "", 1)
        grade = GRADES.get(grade_key)
        context.user_data["grade"] = grade
        await query.edit_message_text(f"✅ تم اختيار: *{grade}*\n\nالآن أرسل في رسالة واحدة:\n*المادة - موضوع الدرس*\nمثال: رياضيات - الضرب", parse_mode="Markdown")
        return

    if data == "mode_pdf_ppt":
        context.user_data["action"] = "pdf_to_ppt"
        await query.edit_message_text("📂 أرسل ملف الـ **PDF** الآن وسأحوله إلى عرض بوربوينت مقسم لشرائح منسقة.", parse_mode="Markdown")
        return

    if data == "mode_pdf_word":
        context.user_data["action"] = "pdf_to_word"
        await query.edit_message_text("📂 أرسل ملف الـ **PDF** الآن وسأقوم بتنظيمه وترتيبه بالذكاء الاصطناعي كملف Word منظم.", parse_mode="Markdown")
        return

    if data == "bot_restart":
        await restart_command(update, context)
        return

    services = {
        "svc_ppt": "📊 بوربوينت 30 شريحة", "svc_exam": "📝 الاختبارات", "svc_remedial": "📈 الخطة العلاجية",
        "svc_portfolio": "🗂 ملف الإنجاز", "svc_performance": "📑 الأداء الوظيفي", "svc_operation": "📅 الخطة التشغيلية",
        "svc_loss": "📚 الفاقد التعليمي", "svc_research": "🎓 بحث جامعي وأكاديمي شامل",
    }

    if data in services:
        context.user_data["current_service"] = data
        context.user_data["service_name"] = services[data]
        grade = context.user_data.get("grade", "")

        if data == "svc_research":
            await query.edit_message_text("🎓 *خدمة البحث الأكاديمي الشامل (جامعة صعدة)*\n\nأرسل الآن **عنوان البحث أو خطة البحث** (مثال: أثر إدارة الأزمات في مشاريع مياه الريف)، وسأقوم بتوليد البحث كاملاً ومفصلاً بجميع الأقسام الأكاديمية.", parse_mode="Markdown")
        else:
            if not grade:
                await query.edit_message_text("🎓 اختر الصف أولاً لتخصيص المحتوى للمرحلة:", reply_markup=grade_menu())
                return
            await query.edit_message_text(f"الخدمة: *{services[data]}*\nالصف: *{grade}*\n\nأرسل الآن: *المادة - الدرس*", parse_mode="Markdown")

async def text_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_text = (update.message.text or "").strip()
    user = update.effective_user
    current_service = context.user_data.get("current_service", "svc_ppt")
    service_name = context.user_data.get("service_name", "عرض بوربوينت")
    grade = context.user_data.get("grade", "دراسات عليا")

    status_msg = await update.message.reply_text("⏳ جارٍ إعداد وتوليد المحتوى المطلوب بدقة واحترافية...")

    try:
        if current_service == "svc_research":
            topic = user_text
            file_name = f"research_{user.id}_{uuid.uuid4().hex[:6]}.docx"
            create_educational_doc_1448(service_code="svc_research", grade="دراسات عليا", subject="بحث أكاديمي", topic=topic, output_path=file_name)

            with open(file_name, "rb") as doc_file:
                await update.message.reply_document(
                    document=doc_file, filename=f"Academic_Research_{topic[:15]}.docx",
                    caption=f"✅ تم إعداد البحث الأكاديمي الشامل والمتكامل بنجاح على غرار نموذج جامعة صعدة\n\n📌 العنوان: {topic}"
                )
            if os.path.exists(file_name): os.remove(file_name)

        else:
            if "-" in user_text:
                subject, topic = user_text.split("-", 1)
            else:
                subject, topic = "عام", user_text

            subject, topic = subject.strip(), topic.strip()
            context.user_data["subject"] = subject

            if current_service == "svc_ppt":
                file_name = f"presentation_{user.id}_{uuid.uuid4().hex[:6]}.pptx"
                create_powerpoint_presentation_full(grade=grade, subject=subject, topic=topic, output_path=file_name)

                with open(file_name, "rb") as ppt_file:
                    await update.message.reply_document(
                        document=ppt_file, filename=f"{subject[:15]}_{topic[:20]}.pptx",
                        caption=f"✅ تم إنشاء العرض التعليمي بتصاميم وألوان متجددة لكل شريحة\n\n🎓 الصف: {grade}\n📚 المادة: {subject}\n📌 الدرس: {topic}"
                    )
                if os.path.exists(file_name): os.remove(file_name)

            else:
                file_name = f"doc_{user.id}_{uuid.uuid4().hex[:6]}.docx"
                create_educational_doc_1448(service_code=current_service, grade=grade, subject=subject, topic=topic, output_path=file_name)

                with open(file_name, "rb") as doc_file:
                    await update.message.reply_document(
                        document=doc_file, filename=f"{service_name[:15]}_{topic[:20]}.docx",
                        caption=f"✅ تم تجهيز مستند Word الاحترافي\n\nالخدمة: {service_name}\nالصف: {grade}\nالمادة: {subject}"
                    )
                if os.path.exists(file_name): os.remove(file_name)

        await status_msg.delete()

    except Exception as exc:
        await status_msg.edit_text(f"⚠️ حدث خطأ أثناء المعالجة. التفاصيل: {str(exc)[:150]}")

async def document_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    doc = update.message.document
    user = update.effective_user
    action = context.user_data.get("action", "pdf_to_word")

    if not doc.file_name.lower().endswith(".pdf"):
        await update.message.reply_text("⚠️ يرجى إرسال ملف بصيغة PDF.")
        return

    status_msg = await update.message.reply_text("⏳ جارٍ استلام ملف الـ PDF ومعالجته...")

    try:
        file = await context.bot.get_file(doc.file_id)
        pdf_bytes = await file.download_as_bytearray()
        
        raw_text = extract_pdf_to_text(bytes(pdf_bytes))

        if action == "pdf_to_ppt":
            output_file = f"ppt_{user.id}_{uuid.uuid4().hex[:6]}.pptx"
            create_ppt_from_pdf_text(raw_text, output_file)
            with open(output_file, "rb") as f:
                await update.message.reply_document(
                    document=f, filename=f"Converted_{doc.file_name[:-4]}.pptx",
                    caption="✅ تم تحويل ملف الـ PDF إلى عرض بوربوينت منسق بنجاح."
                )
        else:
            organized_text = clean_pdf_text_with_ai(raw_text)
            output_file = f"docx_{user.id}_{uuid.uuid4().hex[:6]}.docx"
            
            doc_obj = Document()
            for section in doc_obj.sections:
                section.top_margin, section.bottom_margin = DocxInches(1), DocxInches(1)
                section.left_margin, section.right_margin = DocxInches(1), DocxInches(1)

            title_p = doc_obj.add_paragraph()
            title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run_title = title_p.add_run("ملف مستخرج ومُنظم بدقة")
            run_title.font.size, run_title.font.bold, run_title.font.color.rgb = DocxPt(18), True, DocxRGB(27, 73, 101)

            for block in organized_text.split("\n"):
                clean = block.strip()
                if not clean: continue
                p = doc_obj.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                run = p.add_run(clean)
                run.font.size, run.font.color.rgb = DocxPt(11.5), DocxRGB(30, 41, 59)
                p.paragraph_format.line_spacing, p.paragraph_format.space_after = 1.25, DocxPt(6)

            doc_obj.save(output_file)

            with open(output_file, "rb") as f:
                await update.message.reply_document(
                    document=f, filename=f"Cleaned_{doc.file_name[:-4]}.docx",
                    caption="✅ تم تحويل ملف الـ PDF وتنظيمه في مستند Word احترافي وخالٍ من اللخبطة."
                )

        if os.path.exists(output_file):
            os.remove(output_file)
        await status_msg.delete()
    except Exception as e:
        await status_msg.edit_text(f"⚠️ حدث خطأ أثناء المعالجة: {str(e)[:150]}")

async def handle_ping(request): return web.Response(text="Bot is running perfectly!")

async def start_web_server():
    server = web.Application()
    server.router.add_get("/", handle_ping)
    runner = web.AppRunner(server)
    await runner.setup()
    port = int(os.environ.get("PORT", "10000"))
    site = web.TCPSite(runner, "0.0.0.0", port)
    await site.start()

async def main_async():
    await start_web_server()
    app = Application.builder().token(TELEGRAM_TOKEN).build()
    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("update", restart_command))
    app.add_handler(CommandHandler("restart", restart_command))
    app.add_handler(CallbackQueryHandler(button_handler))
    app.add_handler(MessageHandler(filters.Document.ALL, document_handler))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, text_handler))
    
    print("منصة إنجاز تعمل بكامل الطلبات والمميزات الآن")
    await app.initialize()
    await app.start()
    await app.updater.start_polling()
    while True: await asyncio.sleep(3600)

def main(): asyncio.run(main_async())

if __name__ == "__main__": main()
